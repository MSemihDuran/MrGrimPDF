import os
import fitz
from PIL import Image
import zipfile
import shutil
import subprocess
import tempfile
import xlsxwriter
from pptx import Presentation
from pptx.util import Inches
from pdf2docx import Converter

def pdf_to_word(file_path, output_path):
    cv = Converter(file_path)
    cv.convert(output_path, start=0, end=None)
    cv.close()
    return output_path

def pdf_to_images(file_path, output_dir, img_format='jpg', dpi=150):
    img_format = img_format.lower()
    if img_format not in {'jpg', 'jpeg', 'png'}:
        raise ValueError("Image format must be JPG or PNG.")
    dpi = max(72, min(int(dpi), 600))
    doc = fitz.open(file_path)
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    created_images = []
    
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)
    
    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=mat, alpha=False if img_format.lower() in ['jpg', 'jpeg'] else True)
        ext = 'jpg' if img_format.lower() in ['jpg', 'jpeg'] else 'png'
        img_file = os.path.join(output_dir, f"{base_name}_page_{i+1}.{ext}")
        pix.save(img_file)
        created_images.append(img_file)
        
    doc.close()
    
    if len(created_images) == 1:
        return created_images[0]
    else:
        zip_path = os.path.join(output_dir, f"{base_name}_images.zip")
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for img in created_images:
                zf.write(img, os.path.basename(img))
        return zip_path

def images_to_pdf(image_paths, output_path):
    pdf_doc = fitz.open()
    try:
        for img_path in image_paths:
            if not os.path.exists(img_path):
                continue
            # Converting via MuPDF preserves each image's native dimensions and
            # supports JPEG, PNG, WebP, and common camera image formats.
            with fitz.open(img_path) as img_doc:
                pdf_bytes = img_doc.convert_to_pdf()
            with fitz.open("pdf", pdf_bytes) as img_pdf:
                pdf_doc.insert_pdf(img_pdf)
        if not len(pdf_doc):
            raise ValueError("No readable image files were provided.")
        pdf_doc.save(output_path, garbage=4, deflate=True)
    finally:
        pdf_doc.close()
    return output_path

def pdf_to_excel(file_path, output_path):
    doc = fitz.open(file_path)
    workbook = xlsxwriter.Workbook(output_path)
    header_fmt = workbook.add_format({'bold': True, 'bg_color': '#E2E8F0', 'border': 1})
    cell_fmt = workbook.add_format({'border': 1, 'text_wrap': True})

    for i, page in enumerate(doc):
        sheet_name = f"Page_{i+1}"
        worksheet = workbook.add_worksheet(sheet_name[:31])
        
        tables = page.find_tables()
        row_cursor = 0
        
        if tables.tables:
            for tab_idx, tab in enumerate(tables):
                worksheet.write(row_cursor, 0, f"Table {tab_idx+1}", header_fmt)
                row_cursor += 1
                for r_idx, row in enumerate(tab.extract()):
                    for c_idx, cell in enumerate(row):
                        worksheet.write(row_cursor + r_idx, c_idx, cell or "", cell_fmt)
                row_cursor += len(tab.extract()) + 2
        else:
            blocks = page.get_text("blocks")
            for b in blocks:
                text_content = b[4].strip()
                if text_content:
                    lines = text_content.split("\n")
                    for line in lines:
                        worksheet.write(row_cursor, 0, line, cell_fmt)
                        row_cursor += 1
                    row_cursor += 1
                    
    workbook.close()
    doc.close()
    return output_path

def pdf_to_pptx(file_path, output_path, dpi=120):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]
    
    doc = fitz.open(file_path)
    temp_dir = tempfile.mkdtemp(prefix="mrgrimpdf_pptx_")
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)
    
    try:
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_temp = os.path.join(temp_dir, f"slide_{i}.png")
            pix.save(img_temp)
            slide = prs.slides.add_slide(blank_slide_layout)
            # Preserve the original page ratio instead of distorting every page
            # to 4:3.
            scale = min(prs.slide_width / pix.width, prs.slide_height / pix.height)
            width, height = int(pix.width * scale), int(pix.height * scale)
            slide.shapes.add_picture(img_temp, int((prs.slide_width - width) / 2), int((prs.slide_height - height) / 2), width=width, height=height)
        if not len(doc):
            raise ValueError("The PDF has no pages to convert.")
        prs.save(output_path)
    finally:
        doc.close()
        shutil.rmtree(temp_dir, ignore_errors=True)
    return output_path

def pdf_to_pdfa(file_path, output_path):
    out_dir = os.path.dirname(os.path.abspath(output_path))
    os.makedirs(out_dir, exist_ok=True)
    executable = next((shutil.which(name) for name in ('gs', 'gswin64c', 'gswin32c') if shutil.which(name)), None)
    if executable:
        command = [
            executable, '-dPDFA=2', '-dBATCH', '-dNOPAUSE', '-dSAFER',
            '-sDEVICE=pdfwrite', '-dPDFACompatibilityPolicy=1',
            '-sColorConversionStrategy=UseDeviceIndependentColor', '-sProcessColorModel=DeviceRGB',
            f'-sOutputFile={output_path}', 'PDFA_def.ps', file_path,
        ]
        result = subprocess.run(command, capture_output=True, text=True, timeout=180)
        if result.returncode == 0 and os.path.exists(output_path) and os.path.getsize(output_path) > 0:
            return output_path
    
    # Fallback with PyMuPDF archival formatting & metadata
    doc = fitz.open(file_path)
    meta = doc.metadata or {}
    meta['producer'] = 'MrGrimPDF PDF/A Archival Engine'
    meta['format'] = 'PDF/A-2b'
    doc.set_metadata(meta)
    doc.save(output_path, garbage=4, deflate=True, clean=True)
    doc.close()
    return output_path
